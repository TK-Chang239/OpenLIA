"""Attachment persistence service.

Validates uploads (size, count, mime), writes bytes via
``attachment_storage``, extracts text for non-native-content formats, and
persists ``ChatAttachment`` rows linked to a parent ``ChatMessage``.

See ``planning/specs/systems/composer-attachments-design.md``.
"""

from __future__ import annotations

import io
import uuid
from dataclasses import dataclass
from datetime import UTC, datetime
from typing import Literal

from sqlalchemy.orm import Session

from openlia_server.db.models.content import ChatAttachment
from openlia_server.services import attachment_storage

PER_FILE_MAX_BYTES = 25 * 1024 * 1024
PER_MESSAGE_MAX_FILES = 10

_TEXT_MIMES = frozenset(
    {
        "text/plain",
        "text/markdown",
        "text/csv",
        "text/x-python",
        "text/x-c",
        "text/x-c++",
        "text/x-java",
        "application/json",
    }
)
_IMAGE_MIMES = frozenset({"image/png", "image/jpeg", "image/webp", "image/gif"})
_PDF_MIME = "application/pdf"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_XLSX_MIMES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    }
)
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_ALLOWED_MIMES = _TEXT_MIMES | _IMAGE_MIMES | {_PDF_MIME, _DOCX_MIME, _PPTX_MIME} | _XLSX_MIMES


_ValidationReason = Literal[
    "file_too_large",
    "too_many_files",
    "type_not_allowed",
    "extraction_failed",
]


@dataclass(frozen=True)
class FileUpload:
    filename: str
    mime_type: str
    content: bytes


@dataclass(frozen=True)
class ValidationError:
    filename: str
    reason: _ValidationReason


def validate_uploads(uploads: list[FileUpload]) -> list[ValidationError]:
    """Return per-file errors. Empty list means all uploads are accepted."""
    errors: list[ValidationError] = []
    if len(uploads) > PER_MESSAGE_MAX_FILES:
        errors.append(
            ValidationError(
                filename=f"<{len(uploads)} files>",
                reason="too_many_files",
            )
        )
    for u in uploads:
        if len(u.content) > PER_FILE_MAX_BYTES:
            errors.append(ValidationError(filename=u.filename, reason="file_too_large"))
            continue
        if u.mime_type not in _ALLOWED_MIMES:
            errors.append(ValidationError(filename=u.filename, reason="type_not_allowed"))
    return errors


def persist_attachments(
    db: Session,
    *,
    message_id: str,
    uploads: list[FileUpload],
) -> list[ChatAttachment]:
    """Write each upload to storage, extract text where applicable, and
    persist a ``ChatAttachment`` row. Returns the newly-created rows in the
    same order as ``uploads``.

    Caller is responsible for calling ``validate_uploads`` first; this
    function does not re-validate.
    """
    rows: list[ChatAttachment] = []
    for u in uploads:
        storage_path = attachment_storage.save(u.content, original_filename=u.filename)
        extracted_text, extracted_at = _maybe_extract(u)
        row = ChatAttachment(
            id=str(uuid.uuid4()),
            message_id=message_id,
            filename=u.filename,
            mime_type=u.mime_type,
            size_bytes=len(u.content),
            storage_path=storage_path,
            extracted_text=extracted_text,
            extracted_at=extracted_at,
        )
        db.add(row)
        rows.append(row)
    db.commit()
    for r in rows:
        db.refresh(r)
    return rows


# ─── Text extraction ─────────────────────────────────────────────────────────


def extract_text(upload: FileUpload) -> str | None:
    """Public wrapper: return extracted text for an upload, or None.

    Used by callers (e.g. the v3 engine) that persist attachments
    outside the ``chat_attachments`` table but still want the same
    server-side extraction the chat path uses."""
    return _maybe_extract(upload)[0]


def _maybe_extract(upload: FileUpload) -> tuple[str | None, datetime | None]:
    """Return ``(text, timestamp)`` if extraction is applicable, else
    ``(None, None)``. Images and (intentionally) native PDFs are *not*
    extracted here — the runtime decides at materialization time whether to
    read raw bytes or the cached text. Both are stored to avoid round-tripping
    when the model toggles."""
    mime = upload.mime_type
    extracted_at = datetime.now(UTC)
    try:
        if mime in _TEXT_MIMES:
            return upload.content.decode("utf-8", errors="replace"), extracted_at
        if mime == _PDF_MIME:
            return _extract_pdf(upload.content), extracted_at
        if mime == _DOCX_MIME:
            return _extract_docx(upload.content), extracted_at
        if mime in _XLSX_MIMES:
            return _extract_xlsx(upload.content), extracted_at
        if mime == _PPTX_MIME:
            return _extract_pptx(upload.content), extracted_at
    except Exception:
        # Don't fail the whole upload on a parser error; the runtime
        # materializer will surface a warning if the cache is missing.
        return None, None
    return None, None


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    parts: list[str] = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n\n".join(parts).strip()


def _extract_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            out.append("\t".join("" if v is None else str(v) for v in row))
    return "\n".join(out)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        out.append(f"# Slide {i}")
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                out.append(text)
    return "\n".join(out)
